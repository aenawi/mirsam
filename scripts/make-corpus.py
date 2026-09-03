#!/usr/bin/env python3
"""Generate the golden-corpus decks under tests/fixtures/.

`torture.pptx` and `clean.pptx` (see make-torture-fixture.py) are hand-built
XML in a hand-built container: precise, but nothing like what an application
writes. This script produces decks the shape a real authoring tool produces,
so the corpus exercises the structure mirsam meets in the wild:

  quarterly-report.pptx
      Six slides on python-pptx's default template — a genuine
      PowerPoint-saved theme with a slide master, eleven layouts and the
      English prompt text they carry — with every defect the rule set knows
      spread across placeholders, a text box, a grouped text box, a table and
      speaker notes, next to paragraphs that are correct and paragraphs that
      are English. One paragraph, pasted from a PDF, carries pre-shaped
      presentation forms, so the corpus exercises the one repair that edits
      the letters themselves rather than the properties around them.

  quarterly-report-correct.pptx
      The same deck, authored correctly: explicit direction and alignment on
      every Arabic paragraph, an Arabic language tag, a complex-script font
      wherever a Latin one is set, a native bullet, no controls, no
      presentation forms. mirsam must leave it completely alone. Where
      clean.pptx proves that for a two-paragraph package, this proves it for
      a full template with masters and layouts in the way.

  quarterly-report-impress.pptx   (opt-in: --impress)
      quarterly-report.pptx re-saved by LibreOffice Impress, headless: a
      second application's DrawingML dialect — its own attribute habits, its
      own inheritance flattening — over the same content. Needs an Impress
      installation, not just `soffice`, and is not byte-reproducible across
      machines, since Impress writes what its installed fonts and defaults
      dictate. Not part of the committed corpus until someone with Impress
      runs it and reviews the report it produces.

Deterministic: ZIP timestamps, ordering and compression are normalised after
writing, and core properties are pinned, so re-running reproduces the
committed decks byte for byte.

Usage:
    uv run --with python-pptx scripts/make-corpus.py [--impress]    # or:
    pip install python-pptx && python3 scripts/make-corpus.py

Then regenerate the expected reports:  make golden
"""

from __future__ import annotations

import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FIXTURES = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "tests", "fixtures")

# One timestamp for every entry, so a regenerated deck is the committed deck.
STAMP = (2026, 9, 2, 8, 30, 0)
PINNED = datetime(2026, 9, 2, 8, 30, 0, tzinfo=timezone.utc)

RLM = "‏"
LATIN = "Calibri"
COMPLEX = "Dubai"

# ------------------------------------------------------------------- content


@dataclass
class Para:
    """One paragraph as an author left it.

    Every field is the *defective* state; `correct()` derives what a careful
    author would have written for the same text.
    """

    text: str
    lang: str | None = None  # a:rPr/@lang
    rtl: str | None = None  # a:pPr/@rtl
    algn: str | None = None  # a:pPr/@algn
    latin: str | None = None  # a:latin/@typeface
    cs: str | None = None  # a:cs/@typeface
    level: int = 0
    bullet: bool = False  # a native a:buChar
    clean_text: str | None = None  # what the text should have been

    def is_arabic(self) -> bool:
        return any("ؠ" <= c <= "ۿ" for c in self.text)

    def correct(self) -> Para:
        arabic = self.is_arabic()
        return Para(
            text=self.clean_text if self.clean_text is not None else self.text,
            lang="ar-SA" if arabic else "en-US",
            rtl="1" if arabic else "0",
            algn="r" if arabic else "l",
            latin=self.latin,
            cs=(COMPLEX if (arabic and self.latin) else self.cs),
            level=self.level,
            bullet=True if self.clean_text and self.text.startswith("•") else self.bullet,
        )


@dataclass
class Slide:
    layout: int
    title: Para | None
    body: list[Para]
    notes: Para | None = None
    textbox: list[Para] | None = None
    grouped: list[Para] | None = None
    table: list[list[Para]] | None = None
    name: str = ""


# Layout indexes in python-pptx's default template.
TITLE_SLIDE, TITLE_AND_CONTENT, TITLE_ONLY, BLANK = 0, 1, 5, 6

DECK = [
    # 1. Title slide. The title is pure Arabic declared LTR: the visual order
    #    is the same either way, so the only finding is the language tag.
    #    The subtitle ends in Latin and relies on auto-detection.
    Slide(
        TITLE_SLIDE,
        title=Para("التقرير الفصلي للأداء المؤسسي", lang="en-US", rtl="0"),
        body=[Para("إدارة التخطيط الاستراتيجي | Q4 2026", lang="en-US")],
    ),
    # 2. Bullets, the way they are actually typed: a second level, a run
    #    tagged en-US, a marker typed by hand under a layout that already
    #    supplies one, and the RLM someone added to "fix" a Latin word.
    Slide(
        TITLE_AND_CONTENT,
        title=Para("أبرز النتائج", lang="ar-SA"),
        body=[
            Para("نمت الإيرادات بنسبة 18% مقارنة بالربع السابق", lang="en-US"),
            Para("وارتفع هامش الربح إلى 32%", lang="ar-SA", level=1),
            Para("انخفضت التكاليف التشغيلية بفضل مبادرة ERP الجديدة", lang="ar-SA"),
            Para(
                "• توسّع الفريق إلى 45 موظفاً",
                lang="ar-SA",
                clean_text="توسّع الفريق إلى 45 موظفاً",
            ),
            Para(
                f"تم إطلاق تطبيق Mirsam{RLM} في الإمارات",
                lang="ar-SA",
                clean_text="تم إطلاق تطبيق Mirsam في الإمارات",
            ),
        ],
        notes=Para("ملاحظات المتحدث: ركّز على نمو الإيرادات قبل الانتقال إلى التكاليف.", lang="ar-SA", rtl="1"),
    ),
    # 3. A table. Headers are correct; one cell is hard-aligned left, one is
    #    tagged as English.
    Slide(
        TITLE_ONLY,
        title=Para("مؤشرات الأداء الرئيسية", lang="ar-SA", rtl="1"),
        body=[],
        table=[
            [
                Para("المؤشر", lang="ar-SA", rtl="1"),
                Para("الربع الثالث", lang="ar-SA", rtl="1"),
                Para("الربع الرابع", lang="ar-SA", rtl="1"),
            ],
            [
                Para("الإيرادات (مليون درهم)", lang="ar-SA", rtl="1"),
                Para("120", lang="en-US"),
                Para("142", lang="en-US"),
            ],
            [
                Para("عدد العملاء", lang="en-US", rtl="1"),
                Para("1,850", lang="en-US"),
                Para("2,300", lang="en-US"),
            ],
            [
                Para("رضا العملاء NPS", lang="ar-SA", rtl="1", algn="l"),
                Para("41", lang="en-US"),
                Para("48", lang="en-US"),
            ],
        ],
    ),
    # 4. A text box and a grouped text box. The first paragraph is the worst
    #    case: opens with a Latin acronym, declared LTR, aligned left, Latin
    #    font with an empty complex-script slot, tagged en-US. The second was
    #    pasted from a PDF and carries pre-shaped presentation forms.
    Slide(
        BLANK,
        title=None,
        body=[],
        textbox=[
            Para(
                "GPS يعتمد عليه النظام في تتبّع الشحنات",
                lang="en-US",
                rtl="0",
                algn="l",
                latin=LATIN,
            ),
            Para(
                "الملخص التنفيذي: ﺍﻟﺘﻘﺮﻳﺮ ﺍﻟﻔﺼﻠﻲ",
                lang="ar-SA",
                rtl="1",
                clean_text="الملخص التنفيذي: التقرير الفصلي",
            ),
        ],
        grouped=[
            Para("المصدر: النظام المالي الموحّد", lang="ar-SA", rtl="1", latin="Arial"),
        ],
    ),
    # 5. Authored correctly. Nothing to report.
    Slide(
        TITLE_AND_CONTENT,
        title=Para("التوصيات", lang="ar-SA", rtl="1", algn="r", latin=LATIN, cs=COMPLEX),
        body=[
            Para("الاستثمار في التحوّل الرقمي", lang="ar-SA", rtl="1", algn="r", latin=LATIN, cs=COMPLEX),
            Para("توسيع قاعدة العملاء في دول الخليج", lang="ar-SA", rtl="1", algn="r", latin=LATIN, cs=COMPLEX),
            Para("مراجعة هيكل التكاليف كل ربع", lang="ar-SA", rtl="1", algn="r", latin=LATIN, cs=COMPLEX),
        ],
    ),
    # 6. English. Nothing to report, and nothing to count as Arabic.
    Slide(
        TITLE_AND_CONTENT,
        title=Para("Appendix: Methodology", lang="en-US"),
        body=[
            Para("Figures are unaudited and reported in AED.", lang="en-US"),
            Para("Source: unified finance system, extracted 2026-09-01.", lang="en-US"),
        ],
    ),
]


# ------------------------------------------------------------------- writing


def insert_in_order(parent, child, before: tuple[str, ...]):
    """Insert `child` before the first of `before` present, else append."""
    for tag in before:
        for existing in parent:
            if existing.tag == qn(tag):
                existing.addprevious(child)
                return
    parent.append(child)


def write_para(paragraph, spec: Para):
    """Fill a python-pptx paragraph from a spec, touching only what it names."""
    paragraph.level = spec.level
    pPr = paragraph._p.get_or_add_pPr()
    if spec.algn is not None:
        pPr.set("algn", spec.algn)
    if spec.rtl is not None:
        pPr.set("rtl", spec.rtl)
    if spec.bullet:
        bu = etree.SubElement(pPr, qn("a:buChar"))
        bu.set("char", "•")
        insert_in_order(pPr, bu, ("a:tabLst", "a:defRPr", "a:extLst"))

    run = paragraph.add_run()
    run.text = spec.text
    rPr = run._r.get_or_add_rPr()
    if spec.lang is not None:
        rPr.set("lang", spec.lang)
    rPr.set("dirty", "0")
    if spec.latin is not None:
        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", spec.latin)
        insert_in_order(rPr, latin, ("a:ea", "a:cs", "a:sym", "a:hlinkClick", "a:extLst"))
    if spec.cs is not None:
        cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", spec.cs)
        insert_in_order(rPr, cs, ("a:sym", "a:hlinkClick", "a:extLst"))


def fill_frame(text_frame, paras: list[Para]):
    """Replace a text frame's content with `paras`."""
    first = text_frame.paragraphs[0]
    write_para(first, paras[0])
    for spec in paras[1:]:
        write_para(text_frame.add_paragraph(), spec)


def build(correct: bool) -> Presentation:
    prs = Presentation()
    # An Arabic-first deck declares its reading direction at the top.
    prs.part._element.set("rtl", "1")
    fix = (lambda p: p.correct()) if correct else (lambda p: p)

    for spec in DECK:
        slide = prs.slides.add_slide(prs.slide_layouts[spec.layout])
        placeholders = list(slide.placeholders)

        if spec.title is not None:
            fill_frame(slide.shapes.title.text_frame, [fix(spec.title)])
        if spec.body:
            body = next(p for p in placeholders if p.placeholder_format.idx == 1)
            fill_frame(body.text_frame, [fix(p) for p in spec.body])
        if spec.textbox:
            box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            box.text_frame.word_wrap = True
            fill_frame(box.text_frame, [fix(p) for p in spec.textbox])
        if spec.grouped:
            group = slide.shapes.add_group_shape()
            inner = group.shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(1))
            fill_frame(inner.text_frame, [fix(p) for p in spec.grouped])
        if spec.table:
            rows, cols = len(spec.table), len(spec.table[0])
            frame = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3))
            # A right-to-left table, as an Arabic PowerPoint writes one.
            frame.table._tbl.tblPr.set("rtl", "1")
            for r, row in enumerate(spec.table):
                for c, cell in enumerate(row):
                    fill_frame(frame.table.cell(r, c).text_frame, [fix(cell)])
        if spec.notes is not None:
            fill_frame(slide.notes_slide.notes_text_frame, [fix(spec.notes)])

    core = prs.core_properties
    core.title = "التقرير الفصلي للأداء المؤسسي"
    core.author = "mirsam corpus"
    core.last_modified_by = "mirsam corpus"
    core.created = PINNED
    core.modified = PINNED
    core.revision = 1
    return prs


# ----------------------------------------------------------------- packaging


def normalise(blob: bytes, pin_core: bool = False) -> bytes:
    """Rewrite a package with fixed timestamps and one compression setting.

    Entry order and content are kept. `pin_core` also pins the modified
    timestamp inside docProps/core.xml, which Impress sets to "now".
    """
    out = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(blob)) as src, zipfile.ZipFile(out, "w") as dst:
        for info in src.infolist():
            data = src.read(info)
            if pin_core and info.filename == "docProps/core.xml":
                text = data.decode("utf-8")
                text = re.sub(
                    r"(<dcterms:(?:created|modified)[^>]*>)[^<]*(</dcterms:)",
                    rf"\g<1>{PINNED.strftime('%Y-%m-%dT%H:%M:%SZ')}\g<2>",
                    text,
                )
                data = text.encode("utf-8")
            entry = zipfile.ZipInfo(info.filename, date_time=STAMP)
            entry.compress_type = zipfile.ZIP_DEFLATED
            entry.external_attr = 0o600 << 16
            entry.create_system = 0
            dst.writestr(entry, data, compresslevel=6)
    return out.getvalue()


def save(name: str, blob: bytes):
    path = os.path.join(FIXTURES, name)
    with open(path, "wb") as f:
        f.write(blob)
    with zipfile.ZipFile(path) as z:
        assert z.testzip() is None, f"corrupt entry in {name}"
        n = len(z.infolist())
    print(f"wrote {os.path.relpath(path)}: {n} entries, {len(blob)} bytes")


def serialise(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return normalise(buf.getvalue())


def impress_resave(blob: bytes) -> bytes:
    """Round-trip a deck through LibreOffice Impress, headless."""
    soffice = shutil.which("soffice")
    if soffice is None:
        sys.exit("--impress needs LibreOffice Impress; soffice is not on PATH")
    with tempfile.TemporaryDirectory(prefix="mirsam-corpus-") as tmp:
        src = os.path.join(tmp, "in", "quarterly-report.pptx")
        os.makedirs(os.path.dirname(src))
        with open(src, "wb") as f:
            f.write(blob)
        out_dir = os.path.join(tmp, "out")
        result = subprocess.run(
            [
                soffice,
                "--headless",
                f"-env:UserInstallation=file://{tmp}/profile",
                "--convert-to",
                "pptx",
                "--outdir",
                out_dir,
                src,
            ],
            capture_output=True,
            text=True,
            timeout=300,
        )
        produced = os.path.join(out_dir, "quarterly-report.pptx")
        if result.returncode != 0 or not os.path.exists(produced):
            # `soffice` alone is not enough: without the Impress component it
            # reports "source file could not be loaded" and exits 0.
            sys.exit(
                "Impress could not re-save the deck (is libreoffice-impress installed?)\n"
                + result.stdout
                + result.stderr
            )
        with open(produced, "rb") as f:
            return normalise(f.read(), pin_core=True)


def main():
    os.makedirs(FIXTURES, exist_ok=True)
    defective = serialise(build(correct=False))
    save("quarterly-report.pptx", defective)
    save("quarterly-report-correct.pptx", serialise(build(correct=True)))
    if "--impress" in sys.argv:
        save("quarterly-report-impress.pptx", impress_resave(defective))


if __name__ == "__main__":
    main()
